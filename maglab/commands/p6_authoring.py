"""P6 CLI commands — authoring, comms, gateway, present, hypotheses (§16, §8, §5.10).

Exposes five command surfaces wired via ``register()``:
    write       — F7 journal authoring (Loop C)
    comms       — communications suite (revision, cover-letter, email, abstract, grant)
    gateway     — messaging gateway daemon (setup, start, stop, status, install)
    present     — presentation materials (slides, poster)
    hypotheses  — D1 hypothesis generation

Research integrity contract:
    - Every authoring / comms output carries HUMAN REVIEW REQUIRED.
    - No output is ever auto-sent.  Outputs are written to files for human review.
    - Missing dependency / credentials → clear message + raise typer.Exit(1).

Heavy optional extras (bibtexparser, python-pptx, slack-bolt, discord.py,
python-telegram-bot) are imported LAZILY inside callbacks so that
``maglab --help`` works without [authoring] / [gateway] extras installed.
"""

from __future__ import annotations

import sys
from pathlib import Path
from types import SimpleNamespace
from typing import Annotated

import typer
from rich.console import Console
from rich.panel import Panel
from rich.table import Table

console = Console()

# ---------------------------------------------------------------------------
# Sub-apps
# ---------------------------------------------------------------------------

comms_app = typer.Typer(
    name="comms",
    help="[P6] Academic communications suite (revision, cover-letter, email, abstract, grant).",
    no_args_is_help=True,
)

gateway_app = typer.Typer(
    name="gateway",
    help="[P6] Messaging gateway daemon — Slack / Telegram / Discord (setup, start, stop, status, install).",
    no_args_is_help=True,
)

present_app = typer.Typer(
    name="present",
    help="[P6] Presentation materials (slides, poster).",
    no_args_is_help=True,
)


# ---------------------------------------------------------------------------
# register() — attach all P6 commands to the root maglab app
# ---------------------------------------------------------------------------


def register(app: typer.Typer) -> None:
    """Attach P6 commands to the root maglab app.

    Call this once from the application entry-point after the root ``app``
    has been constructed.  Replaces the dead stubs in ``maglab/cli.py``.

    Parameters
    ----------
    app:
        The root ``typer.Typer`` instance (the ``maglab`` app).
    """
    app.add_typer(comms_app)
    app.add_typer(gateway_app)
    app.add_typer(present_app)
    app.command("write")(write_command)
    app.command("hypotheses")(hypotheses_command)


# ===========================================================================
# write — F7 journal authoring (Loop C)
# ===========================================================================


def write_command(
    results: Annotated[
        str,
        typer.Argument(help="Researcher-provided results summary for the manuscript."),
    ],
    journal: Annotated[
        str,
        typer.Option("--journal", "-j", help="Target journal name (e.g. prl, nature, prb)."),
    ] = "prl",
    loop: Annotated[
        bool,
        typer.Option("--loop/--no-loop", help="Run full Loop C authoring Ralph loop."),
    ] = False,
    dry_run: Annotated[
        bool,
        typer.Option("--dry-run", help="Create output directory structure without LLM calls."),
    ] = False,
    output_dir: Annotated[
        str | None,
        typer.Option("--output-dir", "-o", help="Output directory (default: ./maglab_write/)."),
    ] = None,
) -> None:
    """[P6] Draft a journal manuscript via F7 authoring pipeline (Loop C).

    Results are written to an output directory.  Every output directory
    contains a HUMAN_REVIEW_REQUIRED.txt file.  Auto-submission is
    not supported (§2.4).

    HUMAN REVIEW REQUIRED before any use of these drafts.
    """
    out_dir = Path(output_dir) if output_dir else Path("maglab_write") / journal
    out_dir.mkdir(parents=True, exist_ok=True)

    if dry_run:
        # Create directory skeleton + marker without any LLM calls.
        _write_human_review_marker(out_dir)
        _write_dry_run_manuscript(out_dir / "main.tex", results=results, journal=journal)
        console.print(f"[green]Dry-run complete.[/] Output directory: [bold]{out_dir}[/]")
        console.print(
            "[yellow]HUMAN REVIEW REQUIRED[/] — compile-ready skeleton only. "
            "Replace [FILL] fields and verify all claims before use."
        )
        return

    # Real mode: wire Loop C.
    try:
        from maglab.authoring.bib_manager import BibManager
        from maglab.authoring.data_vault import DataVault
        from maglab.authoring.loop_c import run_loop_c
    except ImportError as exc:
        console.print(
            f"[red]Missing dependency:[/] {exc}\n"
            "Install authoring extras:  pip install maglab[authoring]"
        )
        raise typer.Exit(1) from exc

    vault = DataVault()
    bib = BibManager()

    # Minimal stub LLM callable — requires real LLM credentials at runtime.
    def _llm_fn(system: str, user: str) -> str:  # pragma: no cover
        return _call_llm_or_stub(system, user)

    if loop:
        console.print(
            f"[cyan]Starting Loop C authoring for journal=[bold]{journal}[/]...[/]\n"
            f"[dim]Max 6 iterations | Human sign-off required per section.[/]"
        )
    else:
        console.print(f"[cyan]Drafting manuscript for journal=[bold]{journal}[/]...[/]")

    try:
        result = run_loop_c(
            goal=f"Draft a {journal} manuscript from results: {results[:100]}",
            results_context=results,
            vault=vault,
            bib_manager=bib,
            llm_fn=_llm_fn,
            output_dir=out_dir,
            max_iterations=6 if loop else 1,
            compile_tex=False,
        )
    except Exception as exc:
        console.print(f"[red]Authoring pipeline failed:[/] {exc}")
        raise typer.Exit(1) from exc

    _write_human_review_marker(out_dir)
    console.print(
        f"[green]Draft complete.[/] Sections: {len(result.section_drafts)} | "
        f"Iterations: {result.iterations}"
    )
    console.print(f"  Output: [bold]{out_dir}[/]")
    console.print(
        "[bold yellow]HUMAN REVIEW REQUIRED[/] — do not submit without author review.\n"
        "Per COPE guidelines, AI tools are not listed as authors."
    )


# ===========================================================================
# comms — academic communications suite
# ===========================================================================


def _print_comms_result(result: object, output_path: Path) -> None:
    """Write comms draft to file and print summary."""
    text = getattr(result, "text", str(result))
    wc = getattr(result, "word_count", len(text.split()))
    fills = getattr(result, "fill_markers", [])

    try:
        output_path.write_text(text, encoding="utf-8")
    except OSError as exc:
        console.print(f"[red]Draft write failed:[/] {exc}")
        raise typer.Exit(1) from exc

    console.print(f"[green]Draft saved:[/] [bold]{output_path}[/]")
    console.print(f"  Word count: {wc} | [FILL] markers: {len(fills)}")
    console.print(
        "\n[bold yellow]HUMAN REVIEW REQUIRED[/] — review, personalise [FILL] markers, "
        "then send manually.  This file is NEVER auto-sent."
    )
    if fills:
        console.print("[dim]Fill-in locations:[/]")
        for marker in fills[:5]:
            console.print(f"  • ...{marker}...")


@comms_app.command("revision")
def comms_revision(
    review: Annotated[
        str,
        typer.Option("--review", "-r", help="Path to journal decision letter file."),
    ],
    notes: Annotated[
        str | None,
        typer.Option("--notes", "-n", help="Path to per-comment author notes file."),
    ] = None,
    tone: Annotated[
        str,
        typer.Option("--tone", help="Response tone: formal | respectful | assertive."),
    ] = "formal",
    output: Annotated[
        str | None,
        typer.Option("--output", "-o", help="Output file path (default: revision_letter.txt)."),
    ] = None,
) -> None:
    """[P6] Draft a point-by-point revision response letter.

    HUMAN REVIEW REQUIRED — never auto-sent.
    """
    review_path = Path(review)
    if not review_path.is_file():
        console.print(f"[red]Review file not found:[/] {review!r}")
        raise typer.Exit(1)

    try:
        from maglab.authoring.comms import RevisionLetterAgent
    except ImportError as exc:
        console.print(f"[red]Missing dependency:[/] {exc}")
        raise typer.Exit(1) from exc

    review_text = review_path.read_text(encoding="utf-8")
    notes_text: list[str] = []
    if notes:
        notes_path = Path(notes)
        if notes_path.is_file():
            notes_text = [notes_path.read_text(encoding="utf-8")]
        else:
            console.print(f"[yellow]Notes file not found:[/] {notes!r} — proceeding without notes.")

    def _stub_llm(system: str, user: str) -> str:  # pragma: no cover
        return _call_llm_or_stub(system, user)

    agent = RevisionLetterAgent(llm_fn=_stub_llm)

    with console.status("[dim]Drafting revision letter...[/]"):
        try:
            result = agent.draft(
                {
                    "review_decision": review_text,
                    "comment_notes": notes_text,
                    "tone": tone,
                }
            )
        except Exception as exc:
            console.print(f"[red]Draft failed:[/] {exc}")
            raise typer.Exit(1) from exc

    out = Path(output) if output else Path("revision_letter.txt")
    _print_comms_result(result, out)


@comms_app.command("cover-letter")
def comms_cover_letter(
    journal: Annotated[
        str,
        typer.Option("--journal", "-j", help="Target journal name."),
    ],
    title: Annotated[
        str,
        typer.Option("--title", "-t", help="Manuscript title."),
    ],
    results: Annotated[
        str | None,
        typer.Option("--results", "-r", help="Key results (comma-separated or text)."),
    ] = None,
    output: Annotated[
        str | None,
        typer.Option("--output", "-o", help="Output file path (default: cover_letter.txt)."),
    ] = None,
) -> None:
    """[P6] Draft a cover letter for journal submission (≤250 words).

    HUMAN REVIEW REQUIRED — never auto-sent.
    """
    try:
        from maglab.authoring.comms import CoverLetterAgent
    except ImportError as exc:
        console.print(f"[red]Missing dependency:[/] {exc}")
        raise typer.Exit(1) from exc

    key_results = [r.strip() for r in results.split(",") if r.strip()] if results else []

    def _stub_llm(system: str, user: str) -> str:  # pragma: no cover
        return _call_llm_or_stub(system, user)

    agent = CoverLetterAgent(llm_fn=_stub_llm)

    with console.status("[dim]Drafting cover letter...[/]"):
        try:
            result = agent.draft(
                {
                    "journal": journal,
                    "title": title,
                    "key_results": key_results,
                }
            )
        except Exception as exc:
            console.print(f"[red]Draft failed:[/] {exc}")
            raise typer.Exit(1) from exc

    out = Path(output) if output else Path("cover_letter.txt")
    _print_comms_result(result, out)


@comms_app.command("email")
def comms_email(
    email_type: Annotated[
        str,
        typer.Argument(
            help="Email type: collaboration | question | interview | recommendation | application."
        ),
    ],
    recipient: Annotated[
        str | None,
        typer.Option("--recipient", "-r", help="Recipient name or description."),
    ] = None,
    purpose: Annotated[
        str | None,
        typer.Option("--purpose", "-p", help="Purpose / topic of the email."),
    ] = None,
    output: Annotated[
        str | None,
        typer.Option("--output", "-o", help="Output file path (default: email_draft.txt)."),
    ] = None,
) -> None:
    """[P6] Draft a professional academic email (≤200 words).

    HUMAN REVIEW REQUIRED — never auto-sent.
    """
    valid_types = {"collaboration", "question", "interview", "recommendation", "application"}
    if email_type not in valid_types:
        console.print(
            f"[red]Unknown email type:[/] {email_type!r}. "
            f"Valid types: {', '.join(sorted(valid_types))}"
        )
        raise typer.Exit(1)

    try:
        from maglab.authoring.comms import AcademicEmailAgent
    except ImportError as exc:
        console.print(f"[red]Missing dependency:[/] {exc}")
        raise typer.Exit(1) from exc

    def _stub_llm(system: str, user: str) -> str:  # pragma: no cover
        return _call_llm_or_stub(system, user)

    agent = AcademicEmailAgent(llm_fn=_stub_llm)

    with console.status(f"[dim]Drafting {email_type} email...[/]"):
        try:
            result = agent.draft(
                {
                    "email_type": email_type,
                    "recipient": recipient or "[FILL: recipient name]",
                    "topic": purpose or "[FILL: topic / purpose]",
                }
            )
        except Exception as exc:
            console.print(f"[red]Draft failed:[/] {exc}")
            raise typer.Exit(1) from exc

    out = Path(output) if output else Path("email_draft.txt")
    _print_comms_result(result, out)


@comms_app.command("abstract")
def comms_abstract(
    conference: Annotated[
        str,
        typer.Option("--conference", "-c", help="Conference name (e.g. 'APS March Meeting')."),
    ],
    char_limit: Annotated[
        int,
        typer.Option("--char-limit", help="Character limit for the abstract."),
    ] = 1750,
    results: Annotated[
        str | None,
        typer.Option("--results", "-r", help="Results summary to include."),
    ] = None,
    output: Annotated[
        str | None,
        typer.Option("--output", "-o", help="Output file path (default: abstract_draft.txt)."),
    ] = None,
) -> None:
    """[P6] Draft a conference abstract within a character limit.

    HUMAN REVIEW REQUIRED — never auto-sent.
    """
    try:
        from maglab.authoring.comms import ConferenceAbstractAgent
    except ImportError as exc:
        console.print(f"[red]Missing dependency:[/] {exc}")
        raise typer.Exit(1) from exc

    def _stub_llm(system: str, user: str) -> str:  # pragma: no cover
        return _call_llm_or_stub(system, user)

    agent = ConferenceAbstractAgent(llm_fn=_stub_llm)

    with console.status(f"[dim]Drafting abstract for {conference}...[/]"):
        try:
            result = agent.draft(
                {
                    "conference": conference,
                    "char_limit": char_limit,
                    "results_context": results or "[FILL: describe key results]",
                }
            )
        except Exception as exc:
            console.print(f"[red]Draft failed:[/] {exc}")
            raise typer.Exit(1) from exc

    out = Path(output) if output else Path("abstract_draft.txt")
    _print_comms_result(result, out)

    # Character limit check
    body = getattr(result, "text", "")
    if len(body) > char_limit:
        console.print(
            f"[yellow]Warning:[/] Draft exceeds character limit "
            f"({len(body)} > {char_limit}).  Please trim before submission."
        )


@comms_app.command("grant")
def comms_grant(
    agency: Annotated[
        str,
        typer.Option("--agency", "-a", help="Funding agency (e.g. NSF, DOE, NIH)."),
    ],
    mechanism: Annotated[
        str | None,
        typer.Option("--mechanism", "-m", help="Grant mechanism (e.g. NSF-DMR, DOE-BES)."),
    ] = None,
    aims: Annotated[
        str | None,
        typer.Option("--aims", help="Specific aims or research objectives (text or file path)."),
    ] = None,
    page_limit: Annotated[
        int,
        typer.Option("--page-limit", help="Page limit for specific aims."),
    ] = 2,
    output: Annotated[
        str | None,
        typer.Option("--output", "-o", help="Output file path (default: grant_draft.txt)."),
    ] = None,
) -> None:
    """[P6] Draft grant proposal section text (specific aims, significance, etc.).

    HUMAN REVIEW REQUIRED — never auto-sent.  Budget and co-investigator
    fields are marked [FILL] and must be completed by the researcher.
    """
    try:
        from maglab.authoring.comms import GrantTextAgent
    except ImportError as exc:
        console.print(f"[red]Missing dependency:[/] {exc}")
        raise typer.Exit(1) from exc

    # aims may be a file path
    aims_text = aims or ""
    if aims and Path(aims).is_file():
        aims_text = Path(aims).read_text(encoding="utf-8")

    def _stub_llm(system: str, user: str) -> str:  # pragma: no cover
        return _call_llm_or_stub(system, user)

    agent = GrantTextAgent(llm_fn=_stub_llm)

    with console.status(f"[dim]Drafting grant text for {agency}...[/]"):
        try:
            result = agent.draft(
                {
                    "agency": agency,
                    "mechanism": mechanism or "[FILL: grant mechanism]",
                    "specific_aims": aims_text or "[FILL: specific aims]",
                    "page_limit": page_limit,
                }
            )
        except Exception as exc:
            console.print(f"[red]Draft failed:[/] {exc}")
            raise typer.Exit(1) from exc

    out = Path(output) if output else Path("grant_draft.txt")
    _print_comms_result(result, out)


@comms_app.command("rebuttal")
def comms_rebuttal(
    reviews: Annotated[
        str,
        typer.Option(
            "--reviews",
            "-r",
            help="Path to conference review file, or review text inline.",
        ),
    ],
    notes: Annotated[
        str | None,
        typer.Option("--notes", "-n", help="Path to author rebuttal notes file."),
    ] = None,
    output: Annotated[
        str | None,
        typer.Option("--output", "-o", help="Output file path (default: rebuttal_draft.txt)."),
    ] = None,
) -> None:
    """[P6] Draft a ≤1-page conference rebuttal (≤600 words).

    Clarifies existing results only — no new data introduced.
    HUMAN REVIEW REQUIRED — never auto-sent.
    """
    try:
        from maglab.authoring.comms import RebuttalAgent
    except ImportError as exc:
        console.print(f"[red]Missing dependency:[/] {exc}")
        raise typer.Exit(1) from exc

    # reviews may be a file path or inline text
    reviews_path = Path(reviews)
    if reviews_path.is_file():
        reviews_text: str | list[str] = reviews_path.read_text(encoding="utf-8")
    else:
        reviews_text = reviews  # treat as inline text

    notes_text = ""
    if notes:
        notes_path = Path(notes)
        if notes_path.is_file():
            notes_text = notes_path.read_text(encoding="utf-8")
        else:
            console.print(f"[yellow]Notes file not found:[/] {notes!r} — proceeding without notes.")

    def _stub_llm(system: str, user: str) -> str:  # pragma: no cover
        return _call_llm_or_stub(system, user)

    agent = RebuttalAgent(llm_fn=_stub_llm)

    with console.status("[dim]Drafting conference rebuttal...[/]"):
        try:
            result = agent.draft(
                {
                    "reviews": reviews_text,
                    "author_notes": notes_text or "[FILL: key rebuttal points]",
                }
            )
        except Exception as exc:
            console.print(f"[red]Draft failed:[/] {exc}")
            raise typer.Exit(1) from exc

    out = Path(output) if output else Path("rebuttal_draft.txt")
    _print_comms_result(result, out)


# ===========================================================================
# gateway — messaging gateway daemon
# ===========================================================================


@gateway_app.command("setup")
def gateway_setup(
    config_path: Annotated[
        str | None,
        typer.Option("--config", "-c", help="Path to gateway config YAML/JSON file."),
    ] = None,
) -> None:
    """[P6] Configure the messaging gateway (allowed users, channel IDs, credentials).

    Credentials are stored in a 0600-permission file — never exposed in output.
    """
    import os

    default_cfg = Path.home() / ".maglab" / "gateway.yaml"
    cfg = Path(config_path) if config_path else default_cfg
    cfg.parent.mkdir(parents=True, exist_ok=True)

    if cfg.exists():
        # Check permissions
        mode = oct(cfg.stat().st_mode)[-3:]
        if mode != "600":
            console.print(
                f"[yellow]Warning:[/] Config file {cfg} has permissions {mode} "
                f"(expected 600).  Fix with:  chmod 600 {cfg}"
            )
    else:
        # Write a template config
        template = (
            "# MagLab Gateway Configuration\n"
            "# IMPORTANT: Keep this file 0600 (chmod 600 ~/.maglab/gateway.yaml)\n"
            "\n"
            "slack:\n"
            "  bot_token: ''          # FILL: xoxb-...\n"
            "  app_token: ''          # FILL: xapp-... (Socket Mode)\n"
            "  allowed_users: []      # FILL: list of allowed Slack user IDs\n"
            "  allowed_channels: []   # FILL: list of allowed channel IDs\n"
            "\n"
            "telegram:\n"
            "  bot_token: ''          # FILL: Telegram bot token\n"
            "  allowed_users: []      # FILL: list of allowed chat IDs\n"
            "\n"
            "discord:\n"
            "  bot_token: ''          # FILL: Discord bot token\n"
            "  allowed_users: []      # FILL: list of allowed user IDs\n"
            "  allowed_channels: []   # FILL: list of allowed channel IDs\n"
        )
        cfg.write_text(template, encoding="utf-8")
        os.chmod(cfg, 0o600)

    console.print(f"[green]Gateway config:[/] [bold]{cfg}[/]")
    console.print(
        "  Edit the file to add bot tokens and allowed-user lists.\n"
        "  [yellow]Never share credentials — keep file permissions at 0600.[/]"
    )


@gateway_app.command("start")
def gateway_start(
    foreground: Annotated[
        bool,
        typer.Option("--foreground/--background", help="Run in foreground (default: background)."),
    ] = False,
) -> None:
    """[P6] Start the messaging gateway daemon.

    Background mode forks a subprocess and writes a PID file.
    """
    try:
        from maglab.gateway.runner import _pid_path, is_running
    except ImportError as exc:
        console.print(f"[red]Missing dependency:[/] {exc}")
        raise typer.Exit(1) from exc

    import os

    # MUST read the env var BEFORE calling is_running().
    #
    # When background mode spawns us, the parent writes our PID into the PID
    # file immediately after Popen() returns — well before the child's Python
    # interpreter starts (50-200 ms later).  By the time we reach is_running(),
    # the file already contains our own PID; os.kill(self_pid, 0) always
    # succeeds, so is_running() would return True and we would exit prematurely
    # without ever starting the event loop.
    #
    # The env var tells us the parent has already performed the atomic claim on
    # our behalf.  When it is set we skip is_running() entirely and proceed
    # directly to the event loop.  A direct user invocation (env var absent)
    # still hits the is_running() guard to reject a genuinely-running daemon.
    pid_already_claimed = os.environ.get("MAGLAB_GATEWAY_PID_CLAIMED") == "1"

    if not pid_already_claimed and is_running():
        console.print("[yellow]Gateway is already running.[/] Use 'maglab gateway status'.")
        return

    if foreground:
        # Check whether the background parent already performed the atomic PID
        # claim on our behalf (signalled via MAGLAB_GATEWAY_PID_CLAIMED=1).
        # When the env var is set the child was spawned by background mode; it
        # must NOT attempt a second open("x") on the same file — that would
        # hit FileExistsError because the parent already created it.  Instead
        # the child adopts the file: write_pid() inside _run_gateway_foreground
        # overwrites it with the child's real PID.
        #
        # When the env var is absent the user invoked `gateway start --foreground`
        # directly.  In that case the child IS the first claimer and must perform
        # the atomic exclusive-create to block a concurrent double-start.
        pid_file = _pid_path()
        if not pid_already_claimed:
            # Direct foreground invocation — claim atomically.
            try:
                fd = pid_file.open("x")  # O_CREAT|O_EXCL — fails if file exists
                fd.close()
            except FileExistsError:
                console.print(
                    "[yellow]Gateway is already starting or running.[/] "
                    "Use 'maglab gateway status'."
                )
                return
        # pid_already_claimed=True: parent already holds the lock; we adopt it.
        console.print("[cyan]Starting gateway in foreground mode...[/]")
        _run_gateway_foreground()
    else:
        # Background mode — fork a daemon subprocess.
        #
        # Atomic PID-file claim strategy (preserves Round 5 double-start guard):
        #   1. The PARENT atomically claims the PID file via open("x").
        #      If two concurrent `gateway start` commands race, only the winner
        #      proceeds; the loser gets FileExistsError and exits cleanly.
        #   2. The PARENT spawns `gateway start --foreground` with the env var
        #      MAGLAB_GATEWAY_PID_CLAIMED=1, telling the child that the lock is
        #      already held and it should not re-claim.
        #   3. The PARENT writes proc.pid into the PID file immediately so that
        #      `gateway stop` / `gateway status` can find the real daemon PID.
        #   4. The CHILD's write_pid() call inside _run_gateway_foreground()
        #      then overwrites the file with os.getpid() — the child's own PID —
        #      which is the same value proc.pid refers to.
        import subprocess

        pid_file = _pid_path()
        try:
            fd = pid_file.open("x")  # atomic create — fails if file already exists
            fd.close()
        except FileExistsError:
            console.print(
                "[yellow]Gateway is already starting or running.[/] Use 'maglab gateway status'."
            )
            return

        # Propagate the claim signal to the child so it skips the redundant
        # atomic claim inside the --foreground branch above.
        child_env = {**os.environ, "MAGLAB_GATEWAY_PID_CLAIMED": "1"}

        maglab_exe = sys.executable
        try:
            proc = subprocess.Popen(
                [maglab_exe, "-m", "maglab", "gateway", "start", "--foreground"],
                start_new_session=True,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                env=child_env,
            )
        except Exception:
            # Clean up the sentinel PID file if spawn failed.
            pid_file.unlink(missing_ok=True)
            raise

        # Write the daemon subprocess PID (not the parent CLI's os.getpid()).
        # gateway stop sends SIGTERM to this PID via stop_daemon().
        pid_file.write_text(str(proc.pid))
        console.print(
            f"[green]Gateway started[/] (PID={proc.pid}).  Use 'maglab gateway status' to check."
        )


@gateway_app.command("stop")
def gateway_stop() -> None:
    """[P6] Stop the messaging gateway daemon."""
    try:
        from maglab.gateway.runner import stop_daemon
    except ImportError as exc:
        console.print(f"[red]Missing dependency:[/] {exc}")
        raise typer.Exit(1) from exc

    sent = stop_daemon()
    if sent:
        console.print("[green]Gateway stopped.[/]")
    else:
        console.print("[yellow]Gateway is not running (no PID file found).[/]")


@gateway_app.command("status")
def gateway_status() -> None:
    """[P6] Show the gateway daemon status and PID."""
    try:
        from maglab.gateway.runner import _pid_path, is_running, read_pid
    except ImportError as exc:
        console.print(f"[red]Missing dependency:[/] {exc}")
        raise typer.Exit(1) from exc

    running = is_running()
    pid = read_pid()
    pid_file = _pid_path()

    table = Table(title="Gateway Status", show_lines=False)
    table.add_column("Property", style="cyan")
    table.add_column("Value")
    table.add_row("Status", "[green]Running[/]" if running else "[red]Stopped[/]")
    table.add_row("PID", str(pid) if pid else "(none)")
    table.add_row("PID file", str(pid_file))
    console.print(table)


@gateway_app.command("install")
def gateway_install(
    executable: Annotated[
        str,
        typer.Option("--executable", help="Path to the maglab executable."),
    ] = "maglab",
) -> None:
    """[P6] Install the gateway as a system service (macOS launchd / Linux systemd).

    Requires the gateway config to exist with 0600 permissions.
    """
    try:
        from maglab.gateway.runner import install_service
    except ImportError as exc:
        console.print(f"[red]Missing dependency:[/] {exc}")
        raise typer.Exit(1) from exc

    # Check config permissions before installing
    cfg = Path.home() / ".maglab" / "gateway.yaml"
    if cfg.exists():
        mode = oct(cfg.stat().st_mode)[-3:]
        if mode != "600":
            console.print(
                f"[red]Config file {cfg} has unsafe permissions {mode}.[/]\n"
                "Fix with:  chmod 600 ~/.maglab/gateway.yaml"
            )
            raise typer.Exit(1)

    try:
        service_path = install_service(maglab_executable=executable)
    except (RuntimeError, PermissionError) as exc:
        console.print(f"[red]Service installation failed:[/] {exc}")
        raise typer.Exit(1) from exc

    console.print(f"[green]Service installed:[/] [bold]{service_path}[/]")
    if sys.platform == "darwin":
        console.print("  To activate:  launchctl load " + str(service_path))
    else:
        console.print("  To activate:  systemctl --user enable maglab-gateway")
        console.print("                systemctl --user start maglab-gateway")


# ===========================================================================
# present — slides and poster
# ===========================================================================


@present_app.command("templates")
def present_templates(
    kind: Annotated[
        str,
        typer.Option("--kind", "-k", help="Template kind: all | slides | poster."),
    ] = "all",
    detail: Annotated[
        bool,
        typer.Option("--detail", help="Print constraints, source files, and public references."),
    ] = False,
) -> None:
    """[P6] List installed slide/poster template profiles and how to use them."""
    try:
        from maglab.authoring.present.catalog import list_presentation_templates

        entries = list_presentation_templates(kind)
    except ValueError as exc:
        console.print(f"[red]{exc}[/]")
        raise typer.Exit(2) from exc

    table = Table(title="MagLab presentation templates")
    table.add_column("template", style="bold cyan", no_wrap=True)
    table.add_column("kind")
    table.add_column("formats")
    table.add_column("use case")
    table.add_column("command")

    for entry in entries:
        table.add_row(
            entry.name,
            entry.kind,
            ", ".join(entry.formats),
            entry.use_case,
            entry.command,
        )

    console.print(table)
    if detail:
        for entry in entries:
            console.print(Panel(_template_detail_text(entry), title=entry.name, expand=False))
    console.print()
    console.print("[bold yellow]HUMAN REVIEW REQUIRED[/] for every generated deck or poster.")
    console.print(
        "[dim]Use --template for slides and poster profiles. "
        "Use --format svg, pdf, or beamerposter for posters.[/]"
    )


@present_app.command("slides")
def present_slides(
    results: Annotated[
        str,
        typer.Argument(help="Researcher-provided results summary for the slides."),
    ],
    fmt: Annotated[
        str,
        typer.Option("--format", "-f", help="Slide format: beamer | pptx | marp."),
    ] = "beamer",
    template: Annotated[
        str,
        typer.Option("--template", "-t", help="Template name (e.g. aps-12min, seminar)."),
    ] = "default",
    n_slides: Annotated[
        int | None,
        typer.Option(
            "--n-slides",
            "-n",
            help="Target number of slides. Defaults come from known templates.",
        ),
    ] = None,
    output_dir: Annotated[
        str | None,
        typer.Option("--output-dir", "-o", help="Output directory."),
    ] = None,
    dry_run: Annotated[
        bool,
        typer.Option("--dry-run", help="Create directory structure without LLM calls."),
    ] = False,
) -> None:
    """[P6] Draft presentation slides (beamer PDF / pptx / Marp markdown).

    HUMAN REVIEW REQUIRED — figures and data must be verified before presentation.
    """
    resolved_n_slides = _template_slide_count(template, n_slides)
    out_dir = Path(output_dir) if output_dir else Path("maglab_slides")
    out_dir.mkdir(parents=True, exist_ok=True)

    if dry_run:
        _write_human_review_marker(out_dir)
        _write_dry_run_slides(
            out_dir / f"slides.{_slide_extension(fmt)}",
            results=results,
            fmt=fmt,
            template=template,
            n_slides=resolved_n_slides,
        )
        console.print(f"[green]Dry-run complete.[/] Output: [bold]{out_dir}[/]")
        console.print(
            "[bold yellow]HUMAN REVIEW REQUIRED[/] — format-valid skeleton, "
            "verify before presenting."
        )
        return

    try:
        from maglab.authoring.data_vault import DataVault
        from maglab.authoring.present.slide_drafter import SlideFormat, SlidesDrafter
    except ImportError as exc:
        console.print(f"[red]Missing dependency:[/] {exc}")
        raise typer.Exit(1) from exc

    # Validate format
    try:
        slide_fmt = SlideFormat(fmt)
    except ValueError:
        console.print(f"[red]Unknown format:[/] {fmt!r}.  Valid: beamer, pptx, marp")
        raise typer.Exit(1) from None

    vault = DataVault()

    def _stub_llm(system: str, user: str) -> str:  # pragma: no cover
        return _call_llm_or_stub(system, user)

    drafter = SlidesDrafter(vault=vault, llm_fn=_stub_llm)

    with console.status("[dim]Drafting slides...[/]"):
        try:
            deck = drafter.draft_slides(
                results=results,
                fmt=slide_fmt,
                template=template,
                n_slides=resolved_n_slides,
            )
        except Exception as exc:
            console.print(f"[red]Slide drafting failed:[/] {exc}")
            raise typer.Exit(1) from exc

    # Write output
    ext = _slide_extension(fmt)
    out_file = out_dir / f"slides.{ext}"

    if fmt == "beamer":
        out_file.write_text(deck.to_beamer_tex(), encoding="utf-8")
    elif fmt == "marp":
        out_file.write_text(deck.to_marp_markdown(), encoding="utf-8")
    else:
        # pptx — requires python-pptx (optional extra)
        _write_pptx_deck(deck, out_file)

    _write_human_review_marker(out_dir)

    console.print(f"[green]Slides drafted:[/] [bold]{out_file}[/]  ({len(deck.slides)} slides)")
    console.print(
        "[bold yellow]HUMAN REVIEW REQUIRED[/] — verify all data and figures before presenting."
    )


@present_app.command("poster")
def present_poster(
    results: Annotated[
        str,
        typer.Argument(help="Researcher-provided results summary for the poster."),
    ],
    size: Annotated[
        str | None,
        typer.Option("--size", "-s", help="Poster size (e.g. A0, A1, 96x48in)."),
    ] = None,
    fmt: Annotated[
        str,
        typer.Option("--format", "-f", help="Output format: svg | pdf | beamerposter | tex."),
    ] = "svg",
    title: Annotated[
        str | None,
        typer.Option("--title", "-t", help="Poster title (default: [FILL: poster title])."),
    ] = None,
    template: Annotated[
        str | None,
        typer.Option(
            "--template",
            help="Poster profile: a0-poster | aps-march-poster | beamerposter-a0.",
        ),
    ] = None,
    output_dir: Annotated[
        str | None,
        typer.Option("--output-dir", "-o", help="Output directory."),
    ] = None,
    dry_run: Annotated[
        bool,
        typer.Option("--dry-run", help="Create directory structure without LLM calls."),
    ] = False,
) -> None:
    """[P6] Draft an academic poster layout (SVG / PDF, A0 single-panel).

    Vector layout only — no raster AI image generation (§2.4).
    HUMAN REVIEW REQUIRED before printing or presenting.
    """
    fmt = fmt.strip().lower()
    valid_formats = {"svg", "pdf", "beamerposter", "tex"}
    if fmt not in valid_formats:
        console.print(f"[red]Unknown format:[/] {fmt!r}.  Valid: svg, pdf, beamerposter, tex")
        raise typer.Exit(2)
    resolved_size = _template_poster_size(template, size)

    out_dir = Path(output_dir) if output_dir else Path("maglab_poster")
    out_dir.mkdir(parents=True, exist_ok=True)

    if dry_run:
        _write_human_review_marker(out_dir)
        _write_dry_run_poster(
            out_dir,
            results=results,
            size=resolved_size,
            fmt=fmt,
            template=template,
            title=title or "[FILL: poster title]",
        )
        console.print(f"[green]Dry-run complete.[/] Output: [bold]{out_dir}[/]")
        console.print(
            "[bold yellow]HUMAN REVIEW REQUIRED[/] — format-valid skeleton, verify before printing."
        )
        return

    try:
        from maglab.authoring.data_vault import DataVault
        from maglab.authoring.present.poster_drafter import PosterDrafter
    except ImportError as exc:
        console.print(f"[red]Missing dependency:[/] {exc}")
        raise typer.Exit(1) from exc

    vault = DataVault()

    def _stub_llm(system: str, user: str) -> str:  # pragma: no cover
        return _call_llm_or_stub(system, user)

    drafter = PosterDrafter(vault=vault, llm_fn=_stub_llm)

    with console.status(f"[dim]Drafting {size} poster...[/]"):
        try:
            poster = drafter.draft_poster(
                results=results,
                size=resolved_size,
                fmt=fmt,
                template=template,
                output_dir=out_dir,
                title=title or "[FILL: poster title]",
            )
        except Exception as exc:
            console.print(f"[red]Poster drafting failed:[/] {exc}")
            raise typer.Exit(1) from exc

    _write_human_review_marker(out_dir)

    profile_msg = f", template={template}" if template else ""
    console.print(
        f"[green]Poster drafted:[/] [bold]{poster.path}[/]  "
        f"(format={poster.format}, size={resolved_size}{profile_msg})"
    )
    console.print(
        "[bold yellow]HUMAN REVIEW REQUIRED[/] — verify all data and figures before printing."
    )


# ===========================================================================
# hypotheses — D1 hypothesis generation (§5.10)
# ===========================================================================


def hypotheses_command(
    topic: Annotated[
        str,
        typer.Argument(help="Research topic for hypothesis generation (§5.10 D1)."),
    ],
    n: Annotated[
        int,
        typer.Option("--n", help="Number of hypothesis candidates to generate (1–20)."),
    ] = 5,
    lit_gap: Annotated[
        str | None,
        typer.Option("--lit-gap", help="Literature gap description for grounding."),
    ] = None,
    rng_seed: Annotated[
        int | None,
        typer.Option("--seed", help="RNG seed for reproducible results (testing)."),
    ] = None,
    json_out: Annotated[
        str | None,
        typer.Option("--json-out", help="Write full result as JSON to this path."),
    ] = None,
) -> None:
    """[P6] Generate and rank hypotheses for a research topic (D1, §5.10).

    Uses an Elo tournament (novelty / testability / feasibility / impact) and
    a physics reflection pass.  All outputs are labelled 'AI suggestion' and
    must be verified before use.

    HUMAN REVIEW REQUIRED — hypotheses are AI suggestions, not conclusions.
    """
    try:
        from maglab.core.reasoning import (
            D1HypothesisEngine,
            HypothesisResult,
        )
    except ImportError as exc:
        console.print(f"[red]Missing dependency:[/] {exc}")
        raise typer.Exit(1) from exc

    n = max(1, min(n, 20))

    with console.status(f"[dim]Generating {n} hypothesis candidates for topic: {topic!r}...[/]"):
        try:
            engine = D1HypothesisEngine(n=n, rng_seed=rng_seed)
            result: HypothesisResult = engine.run(
                topic=topic,
                lit_gap=lit_gap or "",
            )
        except Exception as exc:
            console.print(f"[red]Hypothesis generation failed:[/] {exc}")
            raise typer.Exit(1) from exc

    # Render hypothesis cards as Rich Panels
    console.print(
        Panel(
            f"[bold]D1 Hypothesis Engine[/] — Topic: [cyan]{topic}[/]\n[dim]{result.disclaimer}[/]",
            title="[yellow]AI Suggestion — Hypotheses",
            expand=False,
        )
    )

    for rh in result.ranked:
        valid_tag = "[green]Valid[/]" if rh.physical_valid else "[red]PHYSICS ISSUE[/]"
        card_text = (
            f"[bold]#{rh.rank} — {rh.candidate.idea}[/]\n"
            f"  [dim]Novelty:[/] {rh.candidate.novelty_rationale[:120]}\n"
            f"  [dim]Cite-keys:[/] {', '.join(rh.candidate.novelty_cite_keys) or '(none)'}\n"
            f"  [dim]Verify via:[/] {rh.candidate.verification_method or '(not specified)'}\n"
            f"  Feasibility: {rh.candidate.feasibility_score:.2f}  "
            f"Impact: {rh.candidate.impact_score:.2f}  "
            f"Elo: {rh.candidate.elo_rating:.1f}\n"
            f"  Physics: {valid_tag}"
        )
        if not rh.physical_valid and rh.physics_contradiction:
            card_text += f"\n  [red]Warning:[/] {rh.physics_contradiction}"
        console.print(Panel(card_text, title=f"[dim]{rh.ai_label}[/]", expand=False))

    console.print(
        f"\n[dim]Generated {len(result.ranked)} hypothesis card(s).[/]  "
        "[bold yellow]AI suggestions — test before use.[/]"
    )

    # Optional JSON output
    if json_out:
        import json

        out_path = Path(json_out)
        try:
            out_path.write_text(json.dumps(result.to_dict(), indent=2), encoding="utf-8")
            console.print(f"[green]JSON result written:[/] {out_path}")
        except OSError as exc:
            console.print(f"[red]JSON write failed:[/] {exc}")
            raise typer.Exit(1) from exc


# ===========================================================================
# Internal helpers
# ===========================================================================


def _write_human_review_marker(directory: Path) -> None:
    """Write HUMAN_REVIEW_REQUIRED.txt to *directory*."""
    marker = directory / "HUMAN_REVIEW_REQUIRED.txt"
    marker.write_text(
        "HUMAN REVIEW REQUIRED\n\n"
        "This output was drafted with MagLab AI writing assistance.\n"
        "The named authors bear full responsibility for all content, "
        "data, and citations.\n\n"
        "DO NOT SUBMIT OR PRESENT without human review and approval.\n"
        "Per COPE guidelines, AI tools are not listed as authors.\n",
        encoding="utf-8",
    )


def _write_dry_run_manuscript(path: Path, *, results: str, journal: str) -> None:
    """Write a minimal valid manuscript skeleton without any LLM call."""
    summary = _latex_escape_text(results.strip() or "[FILL: verified result summary]")
    path.write_text(
        "\\documentclass[11pt]{article}\n"
        "\\usepackage[margin=1in]{geometry}\n"
        "\\usepackage{graphicx}\n"
        "\\usepackage{siunitx}\n"
        "\\usepackage{hyperref}\n"
        "\\title{[FILL: manuscript title]}\n"
        "\\author{[FILL: author list]}\n"
        "\\date{\\today}\n\n"
        "\\begin{document}\n"
        "\\maketitle\n\n"
        "\\noindent\\textbf{HUMAN REVIEW REQUIRED.} "
        "This is a MagLab dry-run skeleton. Replace all [FILL] fields, "
        "attach verified citations, and bind every number to a DataPoint before use.\n\n"
        "\\begin{abstract}\n"
        "[FILL: journal-specific abstract. Include only verified quantitative claims.]\n"
        "\\end{abstract}\n\n"
        "\\section{Introduction}\n"
        "[FILL: motivation, material system, and research gap.]\n\n"
        "\\section{Methods}\n"
        "[FILL: sample growth, device geometry, measurement protocol, fitting/simulation methods.]\n\n"
        "\\section{Results}\n"
        f"Researcher-provided context for {journal}: {summary}\n\n"
        "\\begin{figure}[t]\n"
        "  \\centering\n"
        "  \\includegraphics[width=0.82\\linewidth]{[FILL: verified figure path]}\n"
        "  \\caption{[FILL: provenance-backed caption with DataPoint IDs.]}\n"
        "\\end{figure}\n\n"
        "\\section{Discussion}\n"
        "[FILL: physical interpretation, limitations, controls, uncertainty.]\n\n"
        "\\section{Data and Code Availability}\n"
        "[FILL: repository, DOI, provenance ledger, and analysis scripts.]\n\n"
        "\\section*{Author Contributions}\n"
        "[FILL: human author contributions. AI tools are not authors.]\n\n"
        "\\bibliographystyle{unsrt}\n"
        "\\bibliography{references}\n"
        "\\end{document}\n",
        encoding="utf-8",
    )


def _write_dry_run_slides(
    path: Path,
    *,
    results: str,
    fmt: str,
    template: str,
    n_slides: int,
) -> None:
    """Write a format-valid slide skeleton without any LLM call."""
    deck = SimpleNamespace(slides=_dry_run_slide_specs(results, n_slides))
    normalized = fmt.strip().lower()
    if normalized == "pptx":
        _write_pptx_deck(deck, path)
        return
    if normalized == "marp":
        path.write_text(_dry_run_marp(deck.slides, template=template), encoding="utf-8")
        return
    if normalized != "beamer":
        console.print(f"[red]Unknown format:[/] {fmt!r}.  Valid: beamer, pptx, marp")
        raise typer.Exit(2)
    path.write_text(_dry_run_beamer(deck.slides, template=template), encoding="utf-8")


def _write_dry_run_poster(
    output_dir: Path,
    *,
    results: str,
    size: str,
    fmt: str,
    template: str | None,
    title: str,
) -> None:
    """Write a format-valid poster skeleton without any LLM call."""
    from maglab.authoring.data_vault import DataVault
    from maglab.authoring.present.poster_drafter import PosterDrafter

    drafter = PosterDrafter(vault=DataVault(), llm_fn=_dry_run_svg_body)
    poster = drafter.draft_poster(
        results=results,
        size=size,
        fmt=fmt,
        template=template,
        output_dir=output_dir,
        title=title,
    )
    if poster.path.suffix.lower() in {".svg", ".tex"}:
        marker = f"MagLab dry-run skeleton - format={fmt}, template={template}, size={size}"
        text = poster.path.read_text(encoding="utf-8")
        comment = (
            f"<!-- {marker} -->\n" if poster.path.suffix.lower() == ".svg" else f"% {marker}\n"
        )
        poster.path.write_text(text.rstrip() + "\n" + comment, encoding="utf-8")


def _dry_run_slide_specs(results: str, n_slides: int) -> list[SimpleNamespace]:
    """Return a concrete, review-ready slide skeleton."""
    summary = results.strip() or "[FILL: verified result summary]"
    base = [
        (
            "[FILL: Title]",
            [
                "HUMAN REVIEW REQUIRED",
                "[FILL: material stack, device, or simulation system]",
                "[FILL: one-sentence physics claim with provenance ID]",
            ],
        ),
        (
            "Research Question",
            [
                "[FILL: bottleneck in magnetism or spintronics workflow]",
                "[FILL: why this measurement/simulation resolves it]",
            ],
        ),
        (
            "Methods and Provenance",
            [
                "[FILL: growth and device geometry]",
                "[FILL: measurement protocol / solver / fitting model]",
                "[FILL: DataPoint IDs and code path]",
            ],
        ),
        ("Verified Result", [summary, "[FILL: uncertainty, units, and source tag]"]),
        (
            "Controls",
            [
                "[FILL: calibration or negative control]",
                "[FILL: symmetry / physics oracle check]",
            ],
        ),
        (
            "Figure Placeholder",
            [
                "[FILL: verified vector figure path]",
                "[FILL: caption with data lineage]",
            ],
        ),
        (
            "Limitations",
            [
                "[FILL: unresolved parameter, model limitation, or sample caveat]",
                "[FILL: next deterministic check]",
            ],
        ),
        ("Takeaway", ["[FILL: one claim only]", "[FILL: next experiment or decision]"]),
    ]
    while len(base) < max(n_slides, 1):
        base.append(
            (
                f"Backup {len(base) + 1}",
                ["[FILL: supporting data]", "[FILL: provenance-backed note]"],
            )
        )
    return [SimpleNamespace(title=title, bullets=bullets) for title, bullets in base[:n_slides]]


def _dry_run_beamer(slides: list[SimpleNamespace], *, template: str) -> str:
    """Return a minimal valid Beamer document for dry-run slides."""
    frames = []
    for slide in slides:
        bullets = "\n".join(f"    \\item {_latex_escape_text(str(item))}" for item in slide.bullets)
        frames.append(
            f"\\begin{{frame}}{{{_latex_escape_text(str(slide.title))}}}\n"
            "  \\begin{itemize}\n"
            f"{bullets}\n"
            "  \\end{itemize}\n"
            "\\end{frame}"
        )
    return (
        "% HUMAN REVIEW REQUIRED\n"
        f"% MagLab dry-run skeleton - format=beamer, template={template}, "
        f"n_slides={len(slides)}\n"
        "\\documentclass[aspectratio=169]{beamer}\n"
        "\\usepackage{graphicx}\n"
        "\\usepackage{siunitx}\n"
        "\\title{[FILL: presentation title]}\n"
        "\\author{[FILL: presenter]}\n"
        f"\\date{{MagLab dry-run skeleton: {template}}}\n\n"
        "\\begin{document}\n"
        "\\maketitle\n\n"
        "\\begin{frame}{Human Review Required}\n"
        "  Replace all [FILL] fields and verify every number, citation, and figure before presenting.\n"
        "\\end{frame}\n\n" + "\n\n".join(frames) + "\n\\end{document}\n"
    )


def _dry_run_marp(slides: list[SimpleNamespace], *, template: str) -> str:
    """Return a minimal valid Marp deck for dry-run slides."""
    lines = [
        "---",
        "marp: true",
        "theme: default",
        'title: "[FILL: presentation title]"',
        "---",
        "",
        "# [FILL: presentation title]",
        "",
        f"MagLab dry-run skeleton: `{template}`",
        "",
        "**HUMAN REVIEW REQUIRED**",
        "",
    ]
    for slide in slides:
        lines.extend(["---", "", f"## {slide.title}", ""])
        lines.extend(f"- {item}" for item in slide.bullets)
        lines.append("")
    return "\n".join(lines)


def _dry_run_svg_body(_system: str, _user: str) -> str:
    """Fallback SVG body used only when a bundled poster template is unavailable."""
    return (
        '<rect x="80" y="80" width="3019" height="4333" fill="#ffffff" stroke="#111111" '
        'stroke-width="8"/>\n'
        '<text x="180" y="220" font-size="82" font-family="Arial" font-weight="700">'
        "[FILL: poster title]</text>\n"
        '<text x="180" y="360" font-size="38" font-family="Arial">HUMAN REVIEW REQUIRED</text>\n'
        '<rect x="180" y="520" width="1280" height="1500" fill="#f7f7f7" stroke="#222"/>\n'
        '<text x="240" y="620" font-size="52" font-family="Arial">Motivation</text>\n'
        '<text x="240" y="720" font-size="34" font-family="Arial">[FILL: research gap]</text>\n'
        '<rect x="1700" y="520" width="1280" height="1500" fill="#f7f7f7" stroke="#222"/>\n'
        '<text x="1760" y="620" font-size="52" font-family="Arial">Verified Results</text>\n'
        '<text x="1760" y="720" font-size="34" font-family="Arial">[FILL: figure path and DataPoint IDs]</text>\n'
        '<rect x="180" y="2300" width="2800" height="900" fill="#f7f7f7" stroke="#222"/>\n'
        '<text x="240" y="2400" font-size="52" font-family="Arial">Takeaway</text>\n'
        '<text x="240" y="2500" font-size="34" font-family="Arial">[FILL: one provenance-backed claim]</text>'
    )


def _latex_escape_text(text: str) -> str:
    """Escape a small amount of user text for LaTeX skeletons."""
    replacements = {
        "\\": r"\textbackslash{}",
        "&": r"\&",
        "%": r"\%",
        "$": r"\$",
        "#": r"\#",
        "_": r"\_",
        "{": r"\{",
        "}": r"\}",
        "~": r"\textasciitilde{}",
        "^": r"\textasciicircum{}",
    }
    return "".join(replacements.get(ch, ch) for ch in text)


def _call_llm_or_stub(system: str, user: str) -> str:  # pragma: no cover
    """Call the default LLM backend, raising a clear error if not configured.

    Imports the LLM backend lazily so that the module can be loaded without
    any [llm] extras installed.
    """
    try:
        from maglab.config import load_config
        from maglab.llm.base import Message, Role
        from maglab.llm.factory import create_llm_backend

        config = load_config()
        backend = create_llm_backend(config)
        result = backend.complete(
            [
                Message(role=Role.SYSTEM, content=system),
                Message(role=Role.USER, content=user),
            ]
        )
        return result.content or ""
    except Exception as exc:  # noqa: BLE001
        raise RuntimeError(
            f"LLM call failed: {exc}\n"
            "Connect a backend with: maglab auth codex | maglab auth <provider> | /connect <provider>"
        ) from exc


def _slide_extension(fmt: str) -> str:
    """Return the file extension for a given slide format."""
    mapping = {"beamer": "tex", "pptx": "pptx", "marp": "md"}
    return mapping.get(fmt, "tex")


def _template_slide_count(template: str, requested: int | None) -> int:
    """Return requested slide count or a source-backed template default."""
    if requested is not None:
        return requested
    normalized = template.strip().lower()
    aliases = {"aps": "aps-12min", "aps-march": "aps-12min", "march-meeting": "aps-12min"}
    target = aliases.get(normalized, normalized)
    return {"aps-12min": 10, "seminar": 24, "internal-update": 8}.get(target, 12)


def _template_poster_size(template: str | None, requested: str | None) -> str:
    """Return requested poster size or a template-specific default."""
    if requested:
        return requested
    normalized = (template or "").strip().lower().replace("_", "-")
    aliases = {"aps": "aps-march-poster", "aps-march": "aps-march-poster"}
    target = aliases.get(normalized, normalized)
    if target in {"aps-march-poster", "march", "march-meeting"}:
        return "96x48in"
    return "A0"


def _poster_extension(fmt: str) -> str:
    """Return the file extension for a given poster format."""
    mapping = {"svg": "svg", "pdf": "pdf", "beamerposter": "tex", "tex": "tex"}
    return mapping.get(fmt, "svg")


def _template_detail_text(entry: object) -> str:
    """Format a presentation template catalog entry."""
    constraints = "\n".join(f"- {item}" for item in getattr(entry, "constraints", ()))
    sources = "\n".join(f"- {path}" for path in getattr(entry, "source_paths", ()))
    refs = "\n".join(f"- {url}" for url in getattr(entry, "reference_urls", ()))
    return (
        f"[bold]Use case[/]\n{getattr(entry, 'use_case', '')}\n\n"
        f"[bold]Constraints[/]\n{constraints or '- none'}\n\n"
        f"[bold]Installed sources[/]\n{sources or '- none'}\n\n"
        f"[bold]Public references[/]\n{refs or '- none'}\n\n"
        f"[bold]Notes[/]\n{getattr(entry, 'notes', '')}"
    )


def _write_pptx_deck(deck: object, out_file: Path) -> None:
    """Write a SlideDeck to a .pptx file using python-pptx (optional dep)."""
    try:
        from pptx import Presentation  # type: ignore[import]
        from pptx.util import Inches  # type: ignore[import]
    except ImportError as exc:
        console.print(
            f"[red]python-pptx not installed:[/] {exc}\nInstall with:  pip install python-pptx"
        )
        raise typer.Exit(1) from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slides_list = getattr(deck, "slides", [])
    for slide_spec in slides_list:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title_ph = slide.shapes.title
        if title_ph:
            title_ph.text = getattr(slide_spec, "title", "Untitled")
        content_ph = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if content_ph:
            tf = content_ph.text_frame
            bullets = getattr(slide_spec, "bullets", [])
            for i, bullet in enumerate(bullets):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = bullet

    prs.save(str(out_file))


def _run_gateway_foreground() -> None:
    """Run the gateway event loop in the foreground (blocking)."""
    import asyncio

    try:
        from maglab.gateway.runner import GatewayRunner, remove_pid, write_pid
    except ImportError as exc:
        console.print(f"[red]Missing gateway dependency:[/] {exc}")
        raise typer.Exit(1) from exc

    runner = GatewayRunner()
    write_pid()

    async def _main() -> None:
        await runner.start()
        console.print("[cyan]Gateway running in foreground.  Press Ctrl+C to stop.[/]")
        try:
            await asyncio.Event().wait()
        except asyncio.CancelledError:
            pass
        finally:
            await runner.stop()
            remove_pid()

    try:
        asyncio.run(_main())
    except KeyboardInterrupt:
        console.print("\n[yellow]Gateway stopped by user.[/]")
    finally:
        remove_pid()
