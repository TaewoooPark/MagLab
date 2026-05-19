"""Slide drafter — beamer / python-pptx / Marp structured deck (§16.6).

Drafts a structured slide deck in the order:
    Title → Motivation → Methods → Results (figure placeholders) → Conclusion.

All numerical values come from DataVault.  Figure placeholders use
``{{figure:SPEC}}`` convention consumed by the figure engine (§12).
"""

from __future__ import annotations

import logging
from collections.abc import Callable
from dataclasses import dataclass, field
from enum import StrEnum
from pathlib import Path

from maglab.authoring.data_vault import DataVault

log = logging.getLogger(__name__)

# Marker written to every output directory
HUMAN_REVIEW_MARKER = (
    "HUMAN REVIEW REQUIRED\n\nAI draft — the presenter bears full responsibility.\n"
)

#: Directory containing bundled presentation templates.
_TEMPLATES_DIR = Path(__file__).parent / "templates"

#: Figure placeholder pattern (consumed by figure engine §12)
FIGURE_PLACEHOLDER = "{{figure:SPEC}}"


class SlideFormat(StrEnum):
    """Supported slide output formats."""

    BEAMER = "beamer"
    PPTX = "pptx"
    MARP = "marp"


@dataclass
class SlideSpec:
    """Specification for a single slide."""

    title: str
    bullets: list[str] = field(default_factory=list)
    figure_placeholder: str | None = None
    notes: str = ""


@dataclass
class SlideDeck:
    """A structured slide deck."""

    slides: list[SlideSpec] = field(default_factory=list)
    format: SlideFormat = SlideFormat.BEAMER
    template_name: str = ""
    human_review_required: bool = True

    @staticmethod
    def _load_template_file(subdir: str, filename: str) -> str | None:
        """Load a bundled presentation template file.

        Looks for ``_TEMPLATES_DIR / subdir / filename``.  Returns the file
        contents as a string, or ``None`` if the file is absent (caller falls
        back to inline generation).
        """
        path = _TEMPLATES_DIR / subdir / filename
        if path.is_file():
            return path.read_text(encoding="utf-8")
        return None

    def to_beamer_tex(self, title: str = "Presentation", author: str = "[FILL: Author]") -> str:
        """Render the deck as a Beamer LaTeX string.

        Attempts to load ``templates/beamer/template.tex`` and substitute
        ``%%TITLE%%``, ``%%AUTHOR%%``, and ``%%SLIDES%%`` tokens.  Falls back
        to inline generation if the template file is absent.
        """
        import datetime

        frames: list[str] = []
        for slide in self.slides:
            frame = f"\\begin{{frame}}{{\\frametitle{{{slide.title}}}}}\n"
            if slide.figure_placeholder:
                frame += f"  % Figure: {slide.figure_placeholder}\n"
                frame += "  \\includegraphics[width=0.8\\textwidth]{[FILL: figure path]}\n"
            if slide.bullets:
                frame += "  \\begin{itemize}\n"
                for b in slide.bullets:
                    frame += f"    \\item {b}\n"
                frame += "  \\end{itemize}\n"
            if slide.notes:
                frame += f"  \\note{{{slide.notes}}}\n"
            frame += "\\end{frame}\n"
            frames.append(frame)
        slides_block = "\n".join(frames)

        template_src = self._load_template_file("beamer", "template.tex")
        if template_src is not None:
            log.debug("SlidesDrafter: using beamer/template.tex")
            today = datetime.date.today().isoformat()
            return (
                template_src.replace("%%TITLE%%", title)
                .replace("%%AUTHOR%%", author)
                .replace("%%SLIDES%%", slides_block)
                .replace("%%DATE%%", today)
            )

        # Inline fallback (original behaviour — backward-compatible)
        preamble = (
            "\\documentclass{beamer}\n"
            "\\usetheme{default}\n"
            "\\usepackage{graphicx}\n"
            "\\usepackage{siunitx}\n"
            f"\\title{{{title}}}\n"
            f"\\author{{{author}}}\n"
            "\\date{\\today}\n\n"
            "\\begin{document}\n"
            "\\maketitle\n\n"
        )
        return preamble + slides_block + "\n\\end{document}\n"

    def to_marp_markdown(self, title: str = "Presentation") -> str:
        """Render the deck as a Marp Markdown string.

        Attempts to load ``templates/marp/template.md`` and substitute
        ``%%TITLE%%`` and ``%%SLIDES%%`` tokens.  Falls back to inline
        generation if the template file is absent.
        """
        import datetime

        slide_lines: list[str] = []
        for slide in self.slides:
            slide_lines.append("---")
            slide_lines.append(f"## {slide.title}")
            slide_lines.append("")
            if slide.figure_placeholder:
                slide_lines.append(f"![{slide.figure_placeholder}]([FILL: figure path])")
                slide_lines.append("")
            for b in slide.bullets:
                slide_lines.append(f"- {b}")
            slide_lines.append("")
        slides_block = "\n".join(slide_lines)

        template_src = self._load_template_file("marp", "template.md")
        if template_src is not None:
            log.debug("SlidesDrafter: using marp/template.md")
            today = datetime.date.today().isoformat()
            return (
                template_src.replace("%%TITLE%%", title)
                .replace("%%SLIDES%%", slides_block)
                .replace("%%DATE%%", today)
            )

        # Inline fallback (original behaviour — backward-compatible)
        lines = [
            "---",
            "marp: true",
            f"title: {title}",
            "---",
            "",
            f"# {title}",
            "",
            "HUMAN REVIEW REQUIRED",
            "",
        ]
        lines.extend(slide_lines)
        return "\n".join(lines)


class SlidesDrafter:
    """Draft a structured slide deck from research results (§16.6).

    Parameters
    ----------
    vault:
        ``DataVault`` for resolving ``{{dp:KEY}}`` placeholders.
    llm_fn:
        LLM callable: ``(system_prompt, user_prompt) → str``.
    """

    _SYSTEM = (
        "You are a presentation design assistant for a magnetism / spintronics researcher.\n\n"
        "RULES:\n"
        "1. Produce a structured slide outline in JSON with keys:\n"
        "   slides: [{title, bullets: [...], figure_placeholder, notes}, ...]\n"
        "2. Use {{dp:KEY}} for numerical values from the data vault.\n"
        "3. Use {{figure:SPEC}} where a figure from the figure engine should appear.\n"
        "4. Do NOT invent numbers or citations.\n"
        "5. Order: Title → Motivation → Methods → Results → Conclusion.\n"
    )

    def __init__(self, vault: DataVault, llm_fn: Callable[[str, str], str]) -> None:
        self._vault = vault
        self._llm = llm_fn

    def draft_slides(
        self,
        results: str,
        fmt: SlideFormat | str = SlideFormat.BEAMER,
        template: str = "default",
        *,
        n_slides: int = 12,
    ) -> SlideDeck:
        """Draft a structured slide deck.

        Parameters
        ----------
        results:
            Researcher-provided results summary.
        fmt:
            Output format (beamer / pptx / marp).
        template:
            Template name hint (e.g. "aps-12min").
        n_slides:
            Target slide count.

        Returns
        -------
        ``SlideDeck``.
        """
        import json

        if isinstance(fmt, str):
            fmt = SlideFormat(fmt)

        user = (
            f"Format: {fmt.value}\n"
            f"Template: {template}\n"
            f"Target slides: {n_slides}\n\n"
            f"Results context:\n{results}\n\n"
            "Produce the slide outline as a JSON object with key 'slides' "
            "containing a list of slide objects."
        )

        raw = self._llm(self._SYSTEM, user)

        # Parse JSON slide specs — fall back to a minimal deck on parse failure.
        slides: list[SlideSpec] = []
        try:
            # Extract JSON block if LLM wrapped it in markdown code fences.
            import re

            json_match = re.search(r"```(?:json)?\s*([\s\S]*?)```", raw)
            json_str = json_match.group(1).strip() if json_match else raw.strip()
            data = json.loads(json_str)
            for s in data.get("slides", []):
                slides.append(
                    SlideSpec(
                        title=s.get("title", "[FILL: slide title]"),
                        bullets=s.get("bullets", []),
                        figure_placeholder=s.get("figure_placeholder"),
                        notes=s.get("notes", ""),
                    )
                )
        except (json.JSONDecodeError, AttributeError):
            log.warning("Slide drafter: JSON parse failed — using fallback single slide.")
            slides = [
                SlideSpec(
                    title="[FILL: slide title]",
                    bullets=["[FILL: content]"],
                    figure_placeholder=FIGURE_PLACEHOLDER,
                )
            ]

        return SlideDeck(slides=slides, format=fmt, template_name=template)

    def export(self, deck: SlideDeck, output_dir: Path, *, title: str = "Presentation") -> Path:
        """Export the deck to the appropriate file format.

        Parameters
        ----------
        deck:
            ``SlideDeck`` to export.
        output_dir:
            Directory to write output files.
        title:
            Presentation title.

        Returns
        -------
        Path to the primary output file.
        """
        output_dir.mkdir(parents=True, exist_ok=True)

        # Write HUMAN_REVIEW_REQUIRED marker
        (output_dir / "HUMAN_REVIEW_REQUIRED.txt").write_text(HUMAN_REVIEW_MARKER, encoding="utf-8")

        if deck.format == SlideFormat.BEAMER:
            tex_content = deck.to_beamer_tex(title=title)
            out_path = output_dir / "slides.tex"
            out_path.write_text(tex_content, encoding="utf-8")
            return out_path

        elif deck.format == SlideFormat.MARP:
            md_content = deck.to_marp_markdown(title=title)
            out_path = output_dir / "slides.md"
            out_path.write_text(md_content, encoding="utf-8")
            return out_path

        elif deck.format == SlideFormat.PPTX:
            return self._export_pptx(deck, output_dir, title)

        else:
            raise ValueError(f"Unknown slide format: {deck.format}")

    def _export_pptx(self, deck: SlideDeck, output_dir: Path, title: str) -> Path:
        """Export as a python-pptx .pptx file."""
        try:
            from pptx import Presentation
        except ImportError:
            log.warning("python-pptx not installed; writing plain text fallback.")
            out_path = output_dir / "slides.txt"
            lines = [HUMAN_REVIEW_MARKER, f"Title: {title}", ""]
            for s in deck.slides:
                lines.append(f"Slide: {s.title}")
                lines.extend(f"  - {b}" for b in s.bullets)
                lines.append("")
            out_path.write_text("\n".join(lines), encoding="utf-8")
            return out_path

        prs = Presentation()
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = "HUMAN REVIEW REQUIRED"

        # Content slides
        for spec in deck.slides:
            layout = prs.slide_layouts[1]  # Title and Content
            sl = prs.slides.add_slide(layout)
            sl.shapes.title.text = spec.title
            tf = sl.placeholders[1].text_frame
            tf.text = ""
            for bullet in spec.bullets:
                para = tf.add_paragraph()
                para.text = bullet
                para.level = 0
            if spec.figure_placeholder:
                tf.add_paragraph().text = f"[Figure: {spec.figure_placeholder}]"

        out_path = output_dir / "slides.pptx"
        prs.save(str(out_path))
        return out_path
